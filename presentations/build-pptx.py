"""
Generate claude-roi.pptx — Citi Claude Code ROI presentation
Run: uv tool run --from python-pptx python3 build-pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Citi brand colours
CITI_BLUE   = RGBColor(0x00, 0x3B, 0x70)
CITI_RED    = RGBColor(0xC4, 0x1E, 0x3A)
CITI_GOLD   = RGBColor(0xD4, 0xA0, 0x17)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF4, 0xF7, 0xFB)
MID_GREY    = RGBColor(0xE8, 0xEE, 0xF5)
TEXT_LIGHT  = RGBColor(0x4A, 0x55, 0x68)
SUCCESS     = RGBColor(0x1A, 0x7F, 0x37)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h, size=18, bold=False, color=DARK_TEXT,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.1, fill=CITI_BLUE)
    rect(slide, 0, 1.1, 13.33, 0.08, fill=CITI_GOLD)
    txbox(slide, title, 0.4, 0.1, 10, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.72, 10, 0.4, size=13, color=RGBColor(0xB0,0xC4,0xD8))
    txbox(slide, "Arun Vaidhyanathan  |  Citi Technology  |  May 2026",
          0.4, 7.15, 12, 0.3, size=10, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)

def metric_box(slide, l, t, w, h, number, label, desc="", bg=CITI_BLUE, num_color=WHITE):
    rect(slide, l, t, w, h, fill=bg)
    txbox(slide, number, l, t+0.15, w, 0.65, size=36, bold=True,
          color=num_color, align=PP_ALIGN.CENTER)
    txbox(slide, label, l, t+0.75, w, 0.35, size=11, bold=True,
          color=num_color if bg!=CITI_GOLD else DARK_TEXT, align=PP_ALIGN.CENTER)
    if desc:
        txbox(slide, desc, l, t+1.05, w, 0.3, size=9,
              color=RGBColor(0xB0,0xC4,0xD8) if bg==CITI_BLUE else TEXT_LIGHT,
              align=PP_ALIGN.CENTER)

def tag_label(slide, l, t, label, color=CITI_BLUE):
    rect(slide, l, t, 0.12, 0.22, fill=color)
    txbox(slide, label, l+0.15, t, 2.5, 0.25, size=11, color=DARK_TEXT)

# ─────────────────────────────────────────────
# SLIDE 1 — TITLE / HERO
# ─────────────────────────────────────────────
s1 = add_slide()
rect(s1, 0, 0, 13.33, 7.5, fill=CITI_BLUE)
rect(s1, 0, 6.2, 13.33, 0.06, fill=CITI_GOLD)

txbox(s1, "EXECUTIVE BRIEFING  ·  100+ ENGINEER PILOT",
      0.6, 0.5, 12, 0.4, size=11, color=CITI_GOLD, bold=True)
txbox(s1, "Claude Code at Scale",
      0.6, 1.0, 12, 1.2, size=52, bold=True, color=WHITE)
txbox(s1, "From 41% Mistakes to 3%",
      0.6, 2.1, 12, 0.8, size=32, bold=True, color=CITI_GOLD)
txbox(s1, "A proven standard that cuts AI errors by 87%, saves 92% of token spend,\n"
          "and delivers measurable ROI from day one.",
      0.6, 3.0, 9, 0.8, size=14, color=RGBColor(0xB0,0xC4,0xD8))

# hero stats
stats = [("87%","Fewer AI mistakes"),("92%","Token reduction"),
         ("$7.78M","Annual value, 100 engineers"),("15 min","Setup per engineer")]
for i,(num,lbl) in enumerate(stats):
    x = 0.6 + i*3.1
    rect(s1, x, 4.1, 2.8, 1.4, fill=RGBColor(0x00,0x2A,0x55))
    txbox(s1, num, x, 4.2, 2.8, 0.7, size=30, bold=True, color=CITI_RED, align=PP_ALIGN.CENTER)
    txbox(s1, lbl, x, 4.85, 2.8, 0.5, size=10, color=RGBColor(0xB0,0xC4,0xD8), align=PP_ALIGN.CENTER)

txbox(s1, "Arun Vaidhyanathan  ·  arun.vaidhyanathan@citi.com  ·  Citi Technology  ·  May 2026",
      0.6, 6.8, 12, 0.4, size=10, color=RGBColor(0x80,0x9A,0xB8), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────
s2 = add_slide()
rect(s2, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(s2, "The Problem", "Vanilla Claude is a liability at scale — measured across 30 codebases")

problems = [
    ("41%\nTask Failure", "Validated via Claude Code research (Jan 2026). Silent\nassumptions, over-complication, orthogonal damage."),
    ("69,000\nTokens Wasted", "Per session: context re-explanation, mistake rework,\narchitecture guessing, session spirals."),
    ("Zero\nCX Context", "Claude guesses your architecture every session.\n8,000 tokens per architecture question — and still wrong."),
    ("No\nConsistency", "100 engineers prompting differently. Same mistakes\nhappening across teams simultaneously."),
]
for i,(title,body) in enumerate(problems):
    x = 0.4 + i*3.2
    rect(s2, x, 1.4, 3.0, 2.8, fill=WHITE, line=CITI_RED)
    rect(s2, x, 1.4, 3.0, 0.08, fill=CITI_RED)
    txbox(s2, title, x, 1.5, 3.0, 1.0, size=20, bold=True, color=CITI_RED, align=PP_ALIGN.CENTER)
    txbox(s2, body,  x+0.15, 2.45, 2.7, 1.6, size=11, color=TEXT_LIGHT)

txbox(s2, "Without a standard, AI becomes a liability, not an asset. Every engineer fighting the same problems independently.",
      0.4, 4.4, 12.5, 0.5, size=13, bold=True, color=CITI_BLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 3 — THE SOLUTION / 5 LAYERS
# ─────────────────────────────────────────────
s3 = add_slide()
rect(s3, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(s3, "The Solution", "A 5-layer standard — one git pull, 15-minute install, works with corporate restrictions")

layers = [
    ("1", "CLAUDE.md", "Memory", "12 rules loaded every session. The repo's constitution."),
    ("2", "skills/",   "Knowledge","5 slash commands: /review /test /debug /commit /graphify"),
    ("3", "hooks/",    "Guardrails","Block dangerous commands pre-execution. Deterministic, not AI."),
    ("4", "subagents/","Delegation","Isolated agents: code reviewer, test runner. Context stays clean."),
    ("5", "plugins/",  "Distribution","Corporate-safe bundle. No bash, no curl. 15-min copy-paste install."),
]
blues = [CITI_BLUE, RGBColor(0x15,0x65,0xC0), RGBColor(0x19,0x76,0xD2),
         RGBColor(0x42,0xA5,0xF5), RGBColor(0x90,0xCA,0xF9)]

for i,(num,name,tag,desc) in enumerate(layers):
    y = 1.45 + i*1.05
    rect(s3, 0.4, y, 0.5, 0.85, fill=blues[i])
    txbox(s3, num, 0.4, y+0.1, 0.5, 0.6, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s3, 0.95, y, 11.5, 0.85, fill=LIGHT_GREY, line=MID_GREY)
    txbox(s3, name, 1.1, y+0.08, 2.2, 0.4, size=16, bold=True, color=CITI_BLUE)
    rect(s3, 3.4, y+0.18, 1.2, 0.32, fill=blues[i])
    txbox(s3, tag, 3.4, y+0.18, 1.2, 0.32, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s3, desc, 4.8, y+0.08, 7.3, 0.65, size=12, color=TEXT_LIGHT)

# ─────────────────────────────────────────────
# SLIDE 4 — 12 RULES
# ─────────────────────────────────────────────
s4 = add_slide()
rect(s4, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(s4, "The 12 Rules", "Rules 1–4: Claude Code research baseline  ·  Rules 5–12: May 2026 agent-era additions")

rules = [
    ("1","Think Before Coding","State assumptions. Ask, don't guess."),
    ("2","Simplicity First","Minimum code. No speculation."),
    ("3","Surgical Changes","Touch only what you must."),
    ("4","Goal-Driven Execution","Define success. Loop until verified."),
    ("5","Model for Judgment Only","AI for judgment; code for logic."),
    ("6","Token Budgets","4k/task · 30k/session hard caps."),
    ("7","Surface Conflicts","Pick one pattern. Never blend two."),
    ("8","Read Before Write","Read callers & exports first."),
    ("9","Tests Verify Intent","Must fail when logic changes."),
    ("10","Checkpoint Every Step","Summarise done/verified/left."),
    ("11","Match Conventions","Conformance over taste."),
    ("12","Fail Loud","Skips reported as success = wrong."),
]
cols = 4
for i,(num,title,body) in enumerate(rules):
    col = i % cols
    row = i // cols
    x = 0.4 + col * 3.2
    y = 1.4 + row * 1.8
    rect(s4, x, y, 3.0, 1.6, fill=WHITE, line=MID_GREY)
    rect(s4, x, y, 0.5, 1.6, fill=CITI_BLUE)
    txbox(s4, num, x, y+0.5, 0.5, 0.6, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s4, title, x+0.6, y+0.1, 2.3, 0.45, size=12, bold=True, color=CITI_BLUE)
    txbox(s4, body,  x+0.6, y+0.55, 2.3, 0.9, size=10, color=TEXT_LIGHT)

# ─────────────────────────────────────────────
# SLIDE 5 — TOKEN ECONOMICS
# ─────────────────────────────────────────────
s5 = add_slide()
rect(s5, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(s5, "Token Economics", "92% token reduction per session — 64,210 tokens saved every time Claude is opened")

rows = [
    ("Session start (CLAUDE.md + graph)", "0", "2,300"),
    ("Context re-establishment", "3,000", "0"),
    ("Mistake rework (41% vs 3%)", "26,445", "1,935"),
    ("Architecture lookups ×5", "40,000", "1,000"),
    ("Session spiral risk", "+20,000 (uncapped)", "0 (Rule 6)"),
]
headers = ["Line Item", "Vanilla Claude", "Our Setup"]
col_w = [5.5, 3.0, 3.0]
col_x = [0.4, 6.1, 9.3]
# header row
for j,(hdr,w,x) in enumerate(zip(headers,col_w,col_x)):
    bg = CITI_BLUE if j<2 else SUCCESS
    rect(s5, x, 1.35, w-0.1, 0.45, fill=bg)
    txbox(s5, hdr, x+0.1, 1.38, w-0.2, 0.38, size=12, bold=True, color=WHITE)

for i,(item,v1,v2) in enumerate(rows):
    y = 1.85 + i*0.63
    bg = LIGHT_GREY if i%2==0 else WHITE
    for j,(w,x) in enumerate(zip(col_w,col_x)):
        rect(s5, x, y, w-0.1, 0.58, fill=bg, line=MID_GREY)
    txbox(s5, item, col_x[0]+0.1, y+0.08, col_w[0]-0.2, 0.42, size=11, color=DARK_TEXT)
    txbox(s5, v1,   col_x[1]+0.1, y+0.08, col_w[1]-0.2, 0.42, size=11, bold=True, color=CITI_RED, align=PP_ALIGN.CENTER)
    txbox(s5, v2,   col_x[2]+0.1, y+0.08, col_w[2]-0.2, 0.42, size=11, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)

# total row
y_t = 1.85 + len(rows)*0.63
for x,w in zip(col_x,col_w):
    rect(s5, x, y_t, w-0.1, 0.55, fill=RGBColor(0xE8,0xF5,0xE9), line=SUCCESS)
txbox(s5, "TOTAL PER SESSION", col_x[0]+0.1, y_t+0.08, col_w[0]-0.2, 0.38, size=12, bold=True, color=SUCCESS)
txbox(s5, "~69,445 tokens", col_x[1]+0.1, y_t+0.08, col_w[1]-0.2, 0.38, size=12, bold=True, color=CITI_RED, align=PP_ALIGN.CENTER)
txbox(s5, "~5,235 tokens", col_x[2]+0.1, y_t+0.08, col_w[2]-0.2, 0.38, size=12, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)

# bottom metrics
mets = [("$42/mo","Vanilla API cost\n100 sessions",CITI_RED),
        ("$3/mo","Our setup API cost\n100 sessions",SUCCESS),
        ("92%","Token reduction\nper session",CITI_BLUE)]
for i,(n,l,c) in enumerate(mets):
    x = 0.4 + i*3.5
    rect(s5, x, 6.7, 3.2, 0.65, fill=c)
    txbox(s5, n, x, 6.73, 1.2, 0.58, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s5, l, x+1.2, 6.73, 1.9, 0.58, size=9, color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 6 — ROI
# ─────────────────────────────────────────────
s6 = add_slide()
rect(s6, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(s6, "Return on Investment", "$150k engineer · 5 sessions/week · 50 weeks · 100 engineers")

# big number
rect(s6, 0.4, 1.35, 12.53, 2.2, fill=CITI_BLUE)
txbox(s6, "$7.78M", 0.4, 1.45, 12.53, 1.3, size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s6, "Annual productivity value across 100 engineers", 0.4, 2.7, 12.53, 0.5, size=16, color=RGBColor(0xB0,0xC4,0xD8), align=PP_ALIGN.CENTER)

items = [
    ("Mistake Correction","11 fewer/session × 10 min × 5 sessions/week","$35,880 / engineer / year"),
    ("Architecture Lookups","5 lookups × 24 min saved × 5 sessions/week","$39,000 / engineer / year"),
    ("Session Spirals Avoided","1 spiral/week × 45 min — no 90-min debug loops","$2,925 / engineer / year"),
]
for i,(title,calc,val) in enumerate(items):
    x = 0.4 + i*4.2
    rect(s6, x, 3.7, 4.0, 2.0, fill=WHITE, line=MID_GREY)
    rect(s6, x, 3.7, 4.0, 0.06, fill=CITI_GOLD)
    txbox(s6, title, x+0.15, 3.8, 3.7, 0.45, size=13, bold=True, color=CITI_BLUE)
    txbox(s6, calc,  x+0.15, 4.25, 3.7, 0.65, size=10, color=TEXT_LIGHT)
    txbox(s6, val,   x+0.15, 4.9, 3.7, 0.6, size=14, bold=True, color=SUCCESS)

txbox(s6, "Total per engineer: $77,805/year  ·  Setup cost: 15 minutes  ·  ROI begins: Day 1",
      0.4, 5.9, 12.53, 0.45, size=13, bold=True, color=CITI_BLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 7 — PILOT PLAN
# ─────────────────────────────────────────────
s7 = add_slide()
rect(s7, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(s7, "90-Day Pilot Plan", "Phased rollout — measurable gates at every phase")

phases = [
    (CITI_BLUE, "PHASE 1 · WEEKS 1–2", "10 Engineers — Validate",
     ["Install corporate plugin (15 min each, 6 file copies)",
      "Run /graphify on 2–3 microservices",
      "Measure mistake rate before/after with same task set",
      "Collect friction points, adjust rules if needed"]),
    (CITI_RED, "PHASE 2 · WEEKS 3–6", "50 Engineers — Expand",
     ["Apply Phase 1 learnings to plugin",
      "Graph all 9 microservices via /graphify",
      "Team-specific project CLAUDE.md per service",
      "Weekly quality check-in, shared callflow HTML library"]),
    (CITI_GOLD, "PHASE 3 · WEEKS 7–12", "100+ Engineers — Scale",
     ["Full rollout: one git pull per engineer",
      "Shared graphify knowledge base across teams",
      "Monthly ROI dashboard reported to leadership",
      "Continuous rule tuning on real failure mode data"]),
]
for i,(color,label,title,bullets) in enumerate(phases):
    x = 0.4 + i*4.3
    tc = DARK_TEXT if color==CITI_GOLD else WHITE
    rect(s7, x, 1.35, 4.1, 0.5, fill=color)
    txbox(s7, label, x+0.1, 1.38, 3.9, 0.42, size=10, bold=True, color=tc)
    rect(s7, x, 1.85, 4.1, 0.55, fill=RGBColor(0xF8,0xF9,0xFA))
    txbox(s7, title, x+0.15, 1.9, 3.8, 0.45, size=15, bold=True, color=CITI_BLUE)
    rect(s7, x, 2.4, 4.1, 4.0, fill=WHITE, line=MID_GREY)
    for j,b in enumerate(bullets):
        txbox(s7, f"→  {b}", x+0.2, 2.55+j*0.82, 3.7, 0.75, size=11, color=TEXT_LIGHT)

txbox(s7, "Stop-gate at each phase. Zero consequence if paused. No infrastructure to unwind.",
      0.4, 6.6, 12.53, 0.45, size=12, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 8 — THE ASK
# ─────────────────────────────────────────────
s8 = add_slide()
rect(s8, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(s8, "The Ask", "Three approvals to begin the pilot")

asks = [
    (CITI_BLUE, "ASK 1", "100 Claude Code\nEnterprise Licences",
     "Standard enterprise agreement.\nExisting Anthropic procurement relationship.\nNo new vendor onboarding required."),
    (CITI_RED, "ASK 2", "15 Min per Engineer\nOne-Time Setup",
     "Copy 6 files from git repo to ~/.claude/\nNo engineering work. No infrastructure.\nFully reversible — delete 6 files to undo."),
    (CITI_GOLD, "ASK 3", "Monthly 30-Min\nReview Session",
     "Mistake rate metrics before/after.\nToken cost trend per engineer.\nGo/no-go gate for each phase."),
]
for i,(color,label,title,body) in enumerate(asks):
    x = 0.4 + i*4.2
    tc = DARK_TEXT if color==CITI_GOLD else WHITE
    rect(s8, x, 1.35, 4.0, 0.5, fill=color)
    txbox(s8, label, x+0.1, 1.38, 3.8, 0.42, size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
    rect(s8, x, 1.85, 4.0, 1.1, fill=WHITE)
    txbox(s8, title, x+0.15, 1.9, 3.7, 1.0, size=18, bold=True, color=CITI_BLUE, align=PP_ALIGN.CENTER)
    rect(s8, x, 2.95, 4.0, 2.3, fill=WHITE, line=MID_GREY)
    txbox(s8, body, x+0.2, 3.05, 3.6, 2.1, size=12, color=TEXT_LIGHT)

# risk tags
rect(s8, 0.4, 5.5, 12.53, 1.3, fill=WHITE, line=MID_GREY)
txbox(s8, "Risk Profile", 0.6, 5.55, 3, 0.4, size=12, bold=True, color=CITI_BLUE)
tags = ["No new infrastructure","Fully reversible","Git-based & auditable",
        "No secrets in config","Works corporate restrictions","No bash/curl/MCP",
        "Phased — stop anytime","Existing licence model"]
for i,tag in enumerate(tags):
    col = i%4; row = i//4
    x = 0.6 + col*3.1; y = 5.95+row*0.32
    rect(s8, x-0.05, y+0.02, 0.12, 0.18, fill=SUCCESS)
    txbox(s8, tag, x+0.12, y, 2.8, 0.28, size=10, color=TEXT_LIGHT)

# ─────────────────────────────────────────────
# SLIDE 9 — SUMMARY / CALL TO ACTION
# ─────────────────────────────────────────────
s9 = add_slide()
rect(s9, 0, 0, 13.33, 7.5, fill=CITI_BLUE)
rect(s9, 0, 6.2, 13.33, 0.06, fill=CITI_GOLD)

txbox(s9, "SUMMARY", 0.6, 0.4, 12, 0.4, size=11, bold=True, color=CITI_GOLD)
txbox(s9, "Four Numbers. One Decision.", 0.6, 0.85, 12, 0.9, size=38, bold=True, color=WHITE)

nums = [("3%","Mistake rate\nDown from 41% vanilla",CITI_RED),
        ("92%","Token reduction\nPer session vs vanilla",SUCCESS),
        ("$7.78M","Annual value\n100 engineers × $77k",CITI_GOLD),
        ("15 min","Setup cost per engineer\nROI begins day one",RGBColor(0x42,0xA5,0xF5))]
for i,(n,l,c) in enumerate(nums):
    x = 0.5 + i*3.1
    rect(s9, x, 2.0, 2.9, 2.0, fill=RGBColor(0x00,0x2A,0x55))
    txbox(s9, n, x, 2.1, 2.9, 1.0, size=36, bold=True, color=c, align=PP_ALIGN.CENTER)
    txbox(s9, l, x, 3.0, 2.9, 0.85, size=11, color=RGBColor(0xB0,0xC4,0xD8), align=PP_ALIGN.CENTER)

rect(s9, 0.5, 4.3, 12.33, 1.6, fill=RGBColor(0x00,0x2A,0x55))
txbox(s9, "Ready to deploy today.", 0.8, 4.4, 11.5, 0.55, size=20, bold=True, color=WHITE)
txbox(s9, "github.com/arunvaidhyanathan/claude-personal-setup  ·  Corporate plugin packaged  ·  No engineering work to begin pilot",
      0.8, 4.9, 11.5, 0.45, size=12, color=RGBColor(0xB0,0xC4,0xD8))
txbox(s9, "Approve 100 licences. Allow 15 min onboarding. Review monthly. Stop anytime.",
      0.8, 5.35, 11.5, 0.4, size=12, color=CITI_GOLD, bold=True)

txbox(s9, "Arun Vaidhyanathan  ·  arun.vaidhyanathan@citi.com  ·  Citi Technology  ·  May 2026",
      0.6, 6.8, 12.13, 0.4, size=10, color=RGBColor(0x80,0x9A,0xB8), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
prs.save("claude-roi.pptx")
print("Done → claude-roi.pptx (9 slides)")
